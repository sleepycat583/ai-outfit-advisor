"""
Generate xiaoyi-competition.pptx from the HTML slide content.
Dark tech theme — faithful to the graphify-dark-graph visual design.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Color Palette ──
BG_DARK  = RGBColor(0x06,0x06,0x0C)
BG_MID   = RGBColor(0x0E,0x10,0x20)
TEXT_1   = RGBColor(0xF0,0xEC,0xE4)
TEXT_2   = RGBColor(0xB0,0xA9,0x9E)
TEXT_3   = RGBColor(0x7A,0x74,0x6C)
ACCENT   = RGBColor(0xE8,0xA8,0x7C)  # warm amber
BLUE     = RGBColor(0x7E,0xB8,0xDA)
GREEN    = RGBColor(0x7E,0xD3,0xA4)
PURPLE   = RGBColor(0xB8,0xA4,0xD6)
ROSE     = RGBColor(0xD4,0xA0,0xB9)
DANGER   = RGBColor(0xE0,0x70,0x70)
GLASS_BG = RGBColor(0x12,0x12,0x1C)
GLASS_BORDER = RGBColor(0x30,0x30,0x38)
WHITE_10 = RGBColor(0x25,0x25,0x2E)

W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # blank layout

# ── Helpers ──
def add_slide():
    return prs.slides.add_slide(blank)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1), radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    if radius is None:
        shape.adjustments[0] = 0.08
    return shape

def add_text_box(slide, left, top, width, height, text="", font_size=Pt(14), color=TEXT_1, bold=False, align=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rich_box(slide, left, top, width, height):
    """Return text_frame for manual paragraph building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_para(tf, text, size=Pt(14), color=TEXT_1, bold=False, align=PP_ALIGN.LEFT, spacing=Pt(6), name='Arial'):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    p.space_after = spacing
    return p

# ── Slide number helper ──
def add_snum(slide, num, total=16):
    add_text_box(slide, Inches(11.6), Inches(0.2), Inches(1.5), Inches(0.4),
                 f"{num:02d} / {total:02d}", Pt(10), TEXT_3)

# ── Eyebrow ──
def add_eyebrow(slide, text, top=Inches(0.4)):
    add_text_box(slide, Inches(0.8), top, Inches(10), Inches(0.4),
                 text, Pt(11), TEXT_3, bold=False, font_name='Arial')

# ── Title ──
def add_title(slide, text, top=Inches(0.85)):
    add_text_box(slide, Inches(0.8), top, Inches(11), Inches(0.8),
                 text, Pt(38), TEXT_1, bold=True, font_name='Arial')

# ── Glass card ──
def glass_card(slide, left, top, w, h, tint=None):
    fill = tint if tint else GLASS_BG
    border = tint if tint else GLASS_BORDER
    return add_rect(slide, left, top, w, h, fill_color=fill, border_color=border)

# ── Tag pill ──
def add_tag(slide, left, top, text, color=ACCENT):
    shape = add_rect(slide, left, top, Inches(1.4), Inches(0.32),
                     fill_color=RGBColor(0x22,0x18,0x10), border_color=RGBColor(0x40,0x30,0x22))
    shape.adjustments[0] = 0.5  # pill shape
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = color
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER

# ══════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 1)
add_text_box(s, Inches(0.8), Inches(2.6), Inches(11), Inches(0.4),
             "2026 未来大学联盟 · 新工科赛道 · 校内选拔赛", Pt(13), TEXT_3)
add_text_box(s, Inches(0.8), Inches(3.1), Inches(11), Inches(1.4),
             "小衣", Pt(88), ACCENT, bold=True)
add_text_box(s, Inches(0.8), Inches(4.3), Inches(11), Inches(0.6),
             "基于大模型与 RAG 的大学生智能穿搭顾问", Pt(22), TEXT_1)
add_text_box(s, Inches(0.8), Inches(5.1), Inches(11), Inches(0.4),
             "未来大学联盟    |    队员：[待填写]    |    2026.05", Pt(13), TEXT_2)
tags = ["LangChain","RAG","Agent","Chroma","Qwen3-max","Streamlit"]
for i,t in enumerate(tags):
    add_tag(s, Inches(0.8 + i*1.55), Inches(5.7), t)

# ══════════════════════════════════════════
# SLIDE 2: AGENDA
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 2)
add_eyebrow(s, "AGENDA")
add_title(s, "汇报大纲")
items = [
    ("01","问题与痛点","大学生穿搭困境与市场现状",ACCENT),
    ("02","解决方案","产品定位与核心功能",BLUE),
    ("03","技术架构","RAG + LangChain + Agent 全链路",PURPLE),
    ("04","核心技术","三大技术亮点深度解析",GREEN),
    ("05","知识库与数据","2026 春夏流行趋势与知识管理",ROSE),
    ("06","创新亮点 & 未来规划","差异化优势与迭代路线图",ACCENT),
]
for i,(num,title,desc,clr) in enumerate(items):
    row = i // 2; col = i % 2
    x = Inches(0.8 + col*6.0); y = Inches(1.9 + row*0.82)
    card = glass_card(s, x, y, Inches(5.7), Inches(0.7))
    tf = card.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = ""
    add_para(tf, num, Pt(24), clr, bold=True, spacing=Pt(0))
    tf.paragraphs[0].text = f"{num}    {title}"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = TEXT_1
    tf.paragraphs[0].font.bold = True
    # Move desc to p2
    add_para(tf, f"     {desc}", Pt(11), TEXT_2)

# ══════════════════════════════════════════
# SLIDE 3: PROBLEM
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 3)
add_eyebrow(s, "Part 01 · 问题定义")
add_title(s, "大学生穿搭的四大困境")
problems = [
    ("01","信息过载，决策瘫痪","小红书、抖音、淘宝……每天接收海量穿搭信息，但真正要出门时反而不知道该穿什么。面对衣柜站 10 分钟是常态。"),
    ("02","场景-天气-风格三重匹配困难","面试要正式、约会要出彩、早八要速战速决。还得查天气、看温度。三个变量同时考虑，现有工具没有一站式解决方案。"),
    ("03","穿搭知识碎片化，缺乏系统指导","洗涤养护、尺码选择、颜色搭配、体型修饰——这些\"穿搭基本功\"分散在各处，没有一个地方能一站式查询。"),
    ("04","通用 AI 不懂\"你\"","ChatGPT 能推荐穿搭，但它不知道你的身高体重、风格偏好、衣柜里有什么。通用方案的推荐千人一面，缺乏个性化。"),
]
for i,(num,title,desc) in enumerate(problems):
    y = Inches(1.85 + i*1.25)
    # Left border
    add_rect(s, Inches(0.8), y, Pt(4), Inches(1.05), fill_color=DANGER, border_color=None, radius=Pt(2))
    card = glass_card(s, Inches(0.95), y, Inches(11.5), Inches(1.05),
                      tint=RGBColor(0x18,0x0E,0x0E))
    tf = card.text_frame; tf.word_wrap = True
    add_para(tf, f"{num}    {title}", Pt(15), TEXT_1, bold=True, spacing=Pt(4))
    add_para(tf, f"     {desc}", Pt(12), TEXT_2)

# ══════════════════════════════════════════
# SLIDE 4: MARKET
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 4)
add_eyebrow(s, "Part 01 · 市场分析")
add_title(s, "目标用户与市场潜力")
stats = [
    ("3,800万+","全国在校大学生",ACCENT),
    ("87%","关注个人形象的大学生",BLUE),
    ("3.2 次/月","人均穿搭类搜索频率",GREEN),
    ("¥560亿","大学生服装市场年规模",PURPLE),
]
for i,(val,label,clr) in enumerate(stats):
    x = Inches(0.8 + i*3.05)
    card = glass_card(s, x, Inches(1.85), Inches(2.85), Inches(1.5))
    tf = card.text_frame; tf.word_wrap = True
    add_para(tf, val, Pt(36), clr, bold=True, align=PP_ALIGN.CENTER, spacing=Pt(6))
    add_para(tf, label, Pt(12), TEXT_2, align=PP_ALIGN.CENTER)

# Insight box
insight = glass_card(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.7),
                     tint=RGBColor(0x1A,0x14,0x0E))
tf = insight.text_frame; tf.word_wrap = True
add_para(tf, "核心发现：大学生群体中，76% 的人表示\"经常为穿什么而困扰\"，但仅 12% 使用过 AI 工具辅助穿搭决策——这是一个巨大的供需缺口。", Pt(14), ACCENT)

# 3 feature cards
feats = [
    ("🎯","高频刚需","每天都要穿衣服，决策频率远高于其他消费场景",BLUE),
    ("📱","天然数字化","Z 世代习惯在线获取信息，对 AI 产品接受度高",GREEN),
    ("🔄","强复购属性","季节更替、场景变化带来持续的使用需求",ACCENT),
]
for i,(icon,title,desc,clr) in enumerate(feats):
    x = Inches(0.8 + i*3.95)
    card = glass_card(s, x, Inches(4.7), Inches(3.7), Inches(1.0))
    tf = card.text_frame; tf.word_wrap = True
    add_para(tf, icon, Pt(22), align=PP_ALIGN.CENTER, spacing=Pt(2))
    add_para(tf, title, Pt(13), TEXT_1, bold=True, align=PP_ALIGN.CENTER, spacing=Pt(2))
    add_para(tf, desc, Pt(10), TEXT_2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 5: SOLUTION
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 5)
add_eyebrow(s, "Part 02 · 解决方案")
add_title(s, "小衣：你的私人穿搭主理人")
add_text_box(s, Inches(0.8), Inches(1.55), Inches(11), Inches(0.4),
             "一个 AI 就够了 —— 她会记住你是谁、了解你的偏好、帮你搜天气、给你专业推荐", Pt(15), TEXT_2)
solutions = [
    ("🗣️","多轮对话","自然语言交互，支持追问与细化",ACCENT),
    ("🌤️","天气感知","Agent 自动联网搜索实时天气",BLUE),
    ("👤","千人千面","性别、风格、体型动态注入 Prompt",GREEN),
    ("📚","专业知识库","洗涤/尺码/颜色/体型四大领域",ACCENT),
    ("⚡","流式输出","逐字渲染，交互体验流畅自然",BLUE),
    ("📤","知识库更新",".txt 一键上传，自动向量化入库",GREEN),
]
for i,(icon,title,desc,clr) in enumerate(solutions):
    row = i // 3; col = i % 3
    x = Inches(0.8 + col*4.0); y = Inches(2.1 + row*1.55)
    card = glass_card(s, x, y, Inches(3.7), Inches(1.35))
    tf = card.text_frame; tf.word_wrap = True
    add_para(tf, icon, Pt(28), align=PP_ALIGN.CENTER, spacing=Pt(2))
    add_para(tf, title, Pt(14), TEXT_1, bold=True, align=PP_ALIGN.CENTER, spacing=Pt(4))
    add_para(tf, desc, Pt(11), TEXT_2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 6: ARCHITECTURE
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 6)
add_eyebrow(s, "Part 03 · 技术架构")
add_title(s, "全栈架构：四层解耦")

layers = [
    ("Streamlit 交互层", ["侧边栏穿搭档案","快捷提问按钮","打字机流式渲染"], BLUE),
    ("LangChain 编排层", ["RunnableWithMessageHistory","AgentExecutor","ConsoleLoggingHandler"], PURPLE),
    ("   └ Tool: DuckDuckGo  |  Tool: Retriever", [], ACCENT),
    ("数据持久层", ["Chroma 向量库","File Chat History"], GREEN),
    ("模型服务层", ["Qwen3-max (LLM)","text-embedding-v4","DashScope API"], ACCENT),
]
y = Inches(1.65)
for i,(layer_name, items, clr) in enumerate(layers):
    if layer_name.startswith("   "):
        # Sub-tool row
        box = add_rect(s, Inches(3.5), y, Inches(9.0), Inches(0.38),
                       fill_color=RGBColor(0x1E,0x16,0x0C), border_color=ACCENT)
        tf = box.text_frame
        add_para(tf, layer_name.strip(), Pt(11), ACCENT, align=PP_ALIGN.CENTER)
        y += Inches(0.1)
    else:
        # Layer header
        box = add_rect(s, Inches(0.8), y, Inches(2.5), Inches(0.55),
                       fill_color=RGBColor(0x10,0x14,0x1A), border_color=clr)
        tf = box.text_frame
        add_para(tf, layer_name, Pt(11), clr, bold=True, align=PP_ALIGN.CENTER)
        # Items
        for j, item in enumerate(items):
            ix = Inches(3.6 + j*3.0)
            item_box = add_rect(s, ix, y, Inches(2.8), Inches(0.55),
                                fill_color=RGBColor(0x10,0x14,0x1A), border_color=clr)
            itf = item_box.text_frame
            add_para(itf, item, Pt(10), TEXT_1, align=PP_ALIGN.CENTER)
        y += Inches(0.65)
    # Arrow
    if i < len(layers)-1:
        add_text_box(s, Inches(6.3), y-Inches(0.08), Inches(1), Inches(0.2),
                     "⬇", Pt(14), ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 7: RAG
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 7)
add_eyebrow(s, "Part 04 · 核心技术一")
add_title(s, "RAG 检索增强生成")
add_text_box(s, Inches(0.8), Inches(1.55), Inches(11), Inches(0.4),
             '不是"让 AI 猜"，而是"让 AI 查完再回答"', Pt(15), TEXT_2)

# Top flow
flow_top = [("📄","知识文档\n.txt 原始文本"), ("✂️","文本分割\nRecursiveCharacter"), ("🔢","向量嵌入\ntext-embedding-v4"), ("🗄️","Chroma 存储\n本地持久化")]
for i,(icon,label) in enumerate(flow_top):
    x = Inches(0.8 + i*3.05)
    box = glass_card(s, x, Inches(2.1), Inches(2.7), Inches(1.0))
    tf = box.text_frame; tf.word_wrap = True
    add_para(tf, icon, Pt(22), align=PP_ALIGN.CENTER, spacing=Pt(2))
    add_para(tf, label, Pt(11), TEXT_1, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(s, Inches(x+2.75), Inches(2.4), Inches(0.3), Inches(0.4), "→", Pt(18), ACCENT)

# Bottom flow
add_text_box(s, Inches(5.8), Inches(3.25), Inches(2), Inches(0.3), "⬇ 用户提问时 ⬇", Pt(12), PURPLE, align=PP_ALIGN.CENTER)
flow_bot = [("❓","用户提问\n自然语言"), ("🔍","语义检索\nChroma 相似度匹配"), ("🧩","上下文拼接\n检索结果注入 Prompt"), ("✨","LLM 生成\nQwen3-max 推理")]
for i,(icon,label) in enumerate(flow_bot):
    x = Inches(0.8 + i*3.05)
    box = glass_card(s, x, Inches(3.6), Inches(2.7), Inches(1.0))
    tf = box.text_frame; tf.word_wrap = True
    add_para(tf, icon, Pt(22), align=PP_ALIGN.CENTER, spacing=Pt(2))
    add_para(tf, label, Pt(11), TEXT_1, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(s, Inches(x+2.75), Inches(3.9), Inches(0.3), Inches(0.4), "→", Pt(18), ACCENT)

# Key design note
note = glass_card(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.7), tint=RGBColor(0x0E,0x18,0x12))
tf = note.text_frame; tf.word_wrap = True
add_para(tf, "关键设计：采用 RecursiveCharacterTextSplitter 按自然段落边界分割文本（chunk_size=500, overlap=50），确保语义完整性。检索到的知识片段直接作为 System Prompt 前缀注入。", Pt(12), GREEN)

# ══════════════════════════════════════════
# SLIDE 8: AGENT
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 8)
add_eyebrow(s, "Part 04 · 核心技术二")
add_title(s, "Agent 联网天气搜索")
add_text_box(s, Inches(0.8), Inches(1.55), Inches(11), Inches(0.4),
             "AI 自主决策：什么时候需要查天气？查什么？怎么用？", Pt(15), TEXT_2)
agent_cards = [
    ("🔧","DuckDuckGo Tool","封装为 LangChain Tool，Agent 自主决定调用时机。无需用户手动输入城市——AI 根据上下文推断搜索范围。",BLUE),
    ("🧠","自主决策链","\"用户问明天面试穿搭\" → Agent 判断：需要查天气 → 自动搜索 → 提取温度/天气 → 融入推荐。",ACCENT),
    ("📊","结构化注入","搜索结果以结构化格式注入 Prompt：温度范围 + 天气状况 + 湿度 + 风力 → 驱动穿搭决策。",GREEN),
]
for i,(icon,title,desc,clr) in enumerate(agent_cards):
    x = Inches(0.8 + i*4.0)
    card = glass_card(s, x, Inches(2.1), Inches(3.7), Inches(1.8))
    tf = card.text_frame; tf.word_wrap = True
    add_para(tf, icon, Pt(28), align=PP_ALIGN.CENTER, spacing=Pt(6))
    add_para(tf, title, Pt(14), clr, bold=True, align=PP_ALIGN.CENTER, spacing=Pt(6))
    add_para(tf, desc, Pt(10), TEXT_2, align=PP_ALIGN.CENTER)

# Code log box
code_box = glass_card(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.3), tint=RGBColor(0x08,0x08,0x10))
tf = code_box.text_frame; tf.word_wrap = True
lines = [
    ("> Entering new AgentExecutor chain...", TEXT_3),
    ("Question: 明天参加社团面试，求推荐穿搭", TEXT_1),
    ("Thought: 用户需要面试穿搭建议，我需要先知道明天的天气情况", BLUE),
    ("Action: duckduckgo_search", ACCENT),
    ('Action Input: "北京 2026年5月19日 天气"', TEXT_2),
    ("Observation: 晴转多云，15°C ~ 26°C，微风", GREEN),
    ("> Finished chain.", TEXT_3),
]
for line,clr in lines:
    add_para(tf, line, Pt(10), clr, spacing=Pt(3), name='Consolas')

# ══════════════════════════════════════════
# SLIDE 9: TERMINAL LIFECYCLE
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 9)
add_eyebrow(s, "Part 04 · 核心技术三")
add_title(s, "底层终端全生命周期管理")
add_text_box(s, Inches(0.8), Inches(1.55), Inches(11), Inches(0.4),
             "从 Agent 启动到链执行完毕，每一毫秒都可追溯", Pt(15), TEXT_2)
term_cards = [
    ("🔍","ConsoleLoggingHandler","自定义 LangChain Callback Handler，捕获 Agent 全生命周期事件：on_llm_start / on_tool_start / on_tool_end / on_chain_end",ACCENT),
    ("🏗️","架构演进：V1.0 → V1.2","V1.0 在 AgentExecutor 内部注册 Callback，导致日志丢失。V1.2 迁移至 chain.stream 根节点 config 注入，确保全链路覆盖。",BLUE),
    ("🛡️","DashScope 代理绕过","入口处设置 NO_PROXY 环境变量白名单，解决校园网代理导致的 DashScope 连接失败问题。",GREEN),
]
for i,(icon,title,desc,clr) in enumerate(term_cards):
    x = Inches(0.8 + i*4.0)
    card = glass_card(s, x, Inches(2.1), Inches(3.7), Inches(1.5))
    tf = card.text_frame; tf.word_wrap = True
    add_para(tf, icon, Pt(24), align=PP_ALIGN.CENTER, spacing=Pt(4))
    add_para(tf, title, Pt(13), clr, bold=True, align=PP_ALIGN.CENTER, spacing=Pt(4))
    add_para(tf, desc, Pt(10), TEXT_2, align=PP_ALIGN.CENTER)

# Compare
add_text_box(s, Inches(0.8), Inches(3.85), Inches(3), Inches(0.3),
             "Callback 注入对比", Pt(13), ACCENT, bold=True)
bad = glass_card(s, Inches(0.8), Inches(4.2), Inches(5.5), Inches(1.2),
                 tint=RGBColor(0x18,0x0E,0x0E))
tf = bad.text_frame; tf.word_wrap = True
add_para(tf, "❌ V1.0 方案", Pt(13), DANGER, bold=True, spacing=Pt(4))
add_para(tf, "在 AgentExecutor 构造函数中传入 Callback，仅覆盖 Agent 执行阶段，chain.stream 外层无法感知内部事件。", Pt(11), TEXT_2)

add_text_box(s, Inches(6.3), Inches(4.6), Inches(0.8), Inches(0.4),
             "VS", Pt(20), ACCENT, bold=True, align=PP_ALIGN.CENTER)

good = glass_card(s, Inches(7.0), Inches(4.2), Inches(5.5), Inches(1.2),
                  tint=RGBColor(0x0E,0x18,0x12))
tf = good.text_frame; tf.word_wrap = True
add_para(tf, "✅ V1.2 方案", Pt(13), GREEN, bold=True, spacing=Pt(4))
add_para(tf, "在 chain.stream() 的 config 参数中注入 Callback，覆盖 RunnableWithMessageHistory → AgentExecutor → Tool 的完整调用链。", Pt(11), TEXT_2)

# ══════════════════════════════════════════
# SLIDE 10: KNOWLEDGE BASE
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 10)
add_eyebrow(s, "Part 05 · 知识库与数据")
add_title(s, "2026 春夏流行趋势 × 知识管理")

# Left column
left = glass_card(s, Inches(0.8), Inches(1.85), Inches(5.7), Inches(3.6))
tf = left.text_frame; tf.word_wrap = True
add_para(tf, "🎨 三大流行风格", Pt(15), TEXT_1, bold=True, spacing=Pt(10))
styles = [("Clean Fit","极简智性风 · 质感优先"),("Y2K","千禧辣妹 · 撞色金属"),("City Walk","城市漫游 · 山系户外")]
for name,desc in styles:
    add_para(tf, f"   {name}    {desc}", Pt(12), TEXT_2, spacing=Pt(6))
add_para(tf, "", Pt(8), spacing=Pt(4))
add_para(tf, "📐 60-30-10 黄金配色法则", Pt(15), TEXT_1, bold=True, spacing=Pt(8))
add_para(tf, "   60% 主色调  |  30% 辅助色  |  10% 点缀色", Pt(12), TEXT_2)

# Right column
right = glass_card(s, Inches(6.8), Inches(1.85), Inches(5.7), Inches(3.6))
tf = right.text_frame; tf.word_wrap = True
add_para(tf, "📚 四大知识库领域", Pt(15), TEXT_1, bold=True, spacing=Pt(10))
domains = ["👔 大厂面试穿搭指南","🧥 衣物洗涤养护知识","📏 尺码推荐与体型修饰","🎨 颜色搭配与流行趋势"]
for d in domains:
    add_para(tf, f"   {d}", Pt(12), TEXT_2, spacing=Pt(6))
add_para(tf, "", Pt(8), spacing=Pt(4))
add_para(tf, "🔄 知识库更新机制", Pt(15), TEXT_1, bold=True, spacing=Pt(8))
add_para(tf, "   上传 .txt → MD5 去重 → 文本分割 → 向量嵌入 → Chroma 持久化 → 即时可检索", Pt(11), TEXT_2)

# Bottom insight
note = glass_card(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5),
                  tint=RGBColor(0x1A,0x14,0x0E))
tf = note.text_frame; tf.word_wrap = True
add_para(tf, "核心优势：知识库与 Agent 天气搜索双向互补——静态知识查库，动态信息联网，覆盖穿搭决策的全部信息需求", Pt(13), ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 11: FEATURE DEMO
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 11)
add_eyebrow(s, "Part 05 · 功能演示")
add_title(s, "核心功能体验流程")
demos = [
    ("👤","Step 1: 建立画像","侧边栏选择性别、风格偏好，输入身高体重 → 动态注入 System Prompt",ACCENT),
    ("💬","Step 2: 自然提问","\"明天参加社团面试求推荐穿搭\" → 一键发送或自由输入",BLUE),
    ("✨","Step 3: 智能生成","Agent 搜索天气 → RAG 检索知识 → LLM 生成专业建议 → 流式渲染",GREEN),
]
for i,(icon,title,desc,clr) in enumerate(demos):
    x = Inches(0.8 + i*4.0)
    card = glass_card(s, x, Inches(1.85), Inches(3.7), Inches(1.5))
    tf = card.text_frame; tf.word_wrap = True
    add_para(tf, icon, Pt(28), align=PP_ALIGN.CENTER, spacing=Pt(4))
    add_para(tf, title, Pt(14), clr, bold=True, align=PP_ALIGN.CENTER, spacing=Pt(4))
    add_para(tf, desc, Pt(10), TEXT_2, align=PP_ALIGN.CENTER)

# Example output
demo = glass_card(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(3.0), tint=RGBColor(0x08,0x08,0x10))
tf = demo.text_frame; tf.word_wrap = True
add_para(tf, "▼ 小衣回答示例（实际输出风格）", Pt(11), TEXT_3, spacing=Pt(6))
add_para(tf, "🎯 场景破冰", Pt(13), ACCENT, bold=True, spacing=Pt(4))
add_para(tf, "明天北京晴转多云，15°C~26°C，微风。社团面试是展现你专业性与个性的关键时刻——既要有仪式感，又不能用力过猛。", Pt(11), TEXT_2, spacing=Pt(6))
add_para(tf, "👔 OOTD 灵感", Pt(13), ACCENT, bold=True, spacing=Pt(4))
add_para(tf, "上装：白色纯棉衬衫（面试永不犯错）→ 配浅灰色 V 领针织背心（增加层次感）\n下装：藏青色直筒西裤（Clean Fit 风，修饰腿型）\n鞋履：白色低帮帆布鞋（干净利落，不过于正式）", Pt(11), TEXT_2, spacing=Pt(6))
add_para(tf, "💡 私藏贴士", Pt(13), ACCENT, bold=True, spacing=Pt(4))
add_para(tf, "面试前一天把衣服挂烫一下，细节决定第一印象。随身带一块手表，不用很贵，但守时是面试的基本素养。", Pt(11), TEXT_2)
add_para(tf, "—— 小衣 · 你的私人穿搭主理人", Pt(11), PURPLE, spacing=Pt(0))

# ══════════════════════════════════════════
# SLIDE 12: INNOVATION
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 12)
add_eyebrow(s, "Part 06 · 创新亮点")
add_title(s, "对比传统方案的四大优势")

# Compare 1
add_text_box(s, Inches(0.8), Inches(1.8), Inches(4), Inches(0.3),
             "传统穿搭 App / 小红书", Pt(13), DANGER, bold=True)
bad1 = glass_card(s, Inches(0.8), Inches(2.1), Inches(4.8), Inches(1.3),
                  tint=RGBColor(0x18,0x0E,0x0E))
tf = bad1.text_frame; tf.word_wrap = True
for line in ["被动浏览，无法主动提问","千人一面，没有个性化","无法感知实时天气","知识分散，不成体系"]:
    add_para(tf, f"• {line}", Pt(12), TEXT_2, spacing=Pt(5))

add_text_box(s, Inches(5.9), Inches(2.5), Inches(1.3), Inches(0.4),
             "VS", Pt(24), ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_text_box(s, Inches(7.4), Inches(1.8), Inches(4), Inches(0.3),
             "小衣 AI 穿搭顾问", Pt(13), GREEN, bold=True)
good1 = glass_card(s, Inches(7.4), Inches(2.1), Inches(4.8), Inches(1.3),
                   tint=RGBColor(0x0E,0x18,0x12))
tf = good1.text_frame; tf.word_wrap = True
for line in ["自然语言多轮对话","用户画像动态注入","Agent 实时天气搜索","RAG 系统知识检索"]:
    add_para(tf, f"✓ {line}", Pt(12), TEXT_2, spacing=Pt(5))

# Compare 2
add_text_box(s, Inches(0.8), Inches(3.7), Inches(4), Inches(0.3),
             "通用 AI（ChatGPT 等）", Pt(13), DANGER, bold=True)
bad2 = glass_card(s, Inches(0.8), Inches(4.0), Inches(4.8), Inches(1.3),
                  tint=RGBColor(0x18,0x0E,0x0E))
tf = bad2.text_frame; tf.word_wrap = True
for line in ["没有穿搭领域知识","不知道你的身材特征","无法获取实时信息","回答风格千篇一律"]:
    add_para(tf, f"• {line}", Pt(12), TEXT_2, spacing=Pt(5))

add_text_box(s, Inches(5.9), Inches(4.4), Inches(1.3), Inches(0.4),
             "VS", Pt(24), ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_text_box(s, Inches(7.4), Inches(3.7), Inches(4), Inches(0.3),
             "小衣 AI 穿搭顾问", Pt(13), GREEN, bold=True)
good2 = glass_card(s, Inches(7.4), Inches(4.0), Inches(4.8), Inches(1.3),
                   tint=RGBColor(0x0E,0x18,0x12))
tf = good2.text_frame; tf.word_wrap = True
for line in ["四大领域专业知识库","身高体重风格动态注入","DuckDuckGo 实时搜索","时尚主理人人设语气"]:
    add_para(tf, f"✓ {line}", Pt(12), TEXT_2, spacing=Pt(5))

# Core moat
note = glass_card(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
                  tint=RGBColor(0x1A,0x14,0x0E))
tf = note.text_frame; tf.word_wrap = True
add_para(tf, "核心壁垒：RAG 知识库 + Agent 联网 + 用户画像 三引擎联动，不是简单套壳，而是从数据层到推理层的完整闭环", Pt(13), ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 13: SCENARIOS
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 13)
add_eyebrow(s, "Part 06 · 应用场景")
add_title(s, "覆盖大学生全场景穿搭需求")
scenes = [
    ("💼","面试/实习","互联网 vs 金融 vs 国企，不同行业的着装潜规则，一键获取"),
    ("💕","约会/联谊","根据约会场景推荐匹配风格的完整 OOTD"),
    ("⏰","早八/日常","5 分钟快速出门方案，学生党的早晨救星"),
    ("🎉","社团活动","百团大战、校园晚会——每个活动有对应穿搭策略"),
    ("🌧️","换季/天气突变","Agent 实时感知温度变化，提醒添衣减衣"),
    ("🛒","购物决策辅助","\"这件衣服配我衣柜里什么好？\"——AI 帮你做搭配"),
]
for i,(icon,title,desc) in enumerate(scenes):
    row = i // 3; col = i % 3
    x = Inches(0.8 + col*4.0); y = Inches(1.85 + row*1.6)
    card = glass_card(s, x, y, Inches(3.7), Inches(1.35))
    tf = card.text_frame; tf.word_wrap = True
    add_para(tf, icon, Pt(24), align=PP_ALIGN.CENTER, spacing=Pt(2))
    add_para(tf, title, Pt(13), TEXT_1, bold=True, align=PP_ALIGN.CENTER, spacing=Pt(3))
    add_para(tf, desc, Pt(10), TEXT_2, align=PP_ALIGN.CENTER)

# Feedback
fb = glass_card(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.5),
                tint=RGBColor(0x0C,0x14,0x1A))
tf = fb.text_frame; tf.word_wrap = True
add_para(tf, '用户反馈："以前面试前刷一小时小红书，现在问小衣 30 秒就搞定了" —— 内测用户 @某 211 高校大三学生', Pt(13), BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 14: ROADMAP
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 14)
add_eyebrow(s, "Part 06 · 未来规划")
add_title(s, "功能迭代路线图")

# Timeline
phases = [
    ("V1.2","当前","Agent 天气搜索\nCallback 架构升级\n2026 春夏知识库", True),
    ("V2.0","Q3 2026","多模态输入\n上传照片分析穿搭\n虚拟试穿预览", True),
    ("V2.5","Q4 2026","衣橱数字化\n拍照录入已有衣物\n基于衣橱智能搭配", True),
    ("V3.0","2027","社区化 + 电商\n用户穿搭分享\n一键购买推荐单品", False),
]
# Timeline line
add_rect(s, Inches(1.2), Inches(2.15), Inches(11.0), Pt(1), fill_color=WHITE_10, border_color=None)
for i,(ver,date,desc,active) in enumerate(phases):
    x = Inches(1.2 + i*2.85)
    # Dot
    dot_color = ACCENT if active else WHITE_10
    add_rect(s, x+Inches(1.1), Inches(2.05), Inches(0.22), Inches(0.22),
             fill_color=dot_color, border_color=None, radius=Pt(1))
    add_text_box(s, x, Inches(2.35), Inches(2.5), Inches(0.25),
                 f"{ver}  {date}", Pt(11), ACCENT if active else TEXT_3, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s, x, Inches(2.65), Inches(2.5), Inches(1.0),
                 desc, Pt(11), TEXT_2 if active else TEXT_3, align=PP_ALIGN.CENTER)

# 3 future cards
future_cards = [
    ("🖼️","多模态 V2.0","支持上传穿搭照片，AI 分析并给出改进建议。集成虚拟试穿能力。",BLUE),
    ("👗","衣橱数字化 V2.5","拍照识别已有衣物，构建个人数字衣橱。AI 基于实际衣橱做搭配。",ACCENT),
    ("🌐","社区 + 电商 V3.0","穿搭方案一键分享，构建大学生穿搭社区。智能推荐购买链接。",GREEN),
]
for i,(icon,title,desc,clr) in enumerate(future_cards):
    x = Inches(0.8 + i*4.0)
    card = glass_card(s, x, Inches(4.0), Inches(3.7), Inches(1.7))
    tf = card.text_frame; tf.word_wrap = True
    add_para(tf, icon, Pt(24), align=PP_ALIGN.CENTER, spacing=Pt(4))
    add_para(tf, title, Pt(14), clr, bold=True, align=PP_ALIGN.CENTER, spacing=Pt(4))
    add_para(tf, desc, Pt(11), TEXT_2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 15: TEAM
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 15)
add_eyebrow(s, "团队介绍")
add_title(s, "关于我们")
team = [
    ("👨‍💻","[队员一]","项目负责人 · 全栈开发","负责 LangChain Agent 架构设计与 RAG 系统搭建"),
    ("👩‍🎨","[队员二]","UI/UX · 前端开发","负责 Streamlit 前端交互设计与打字机流式渲染"),
    ("👩‍🔬","[队员三]","知识库 · 数据工程","负责穿搭知识库构建、流行趋势分析与数据标注"),
]
for i,(avatar,name,role,desc) in enumerate(team):
    x = Inches(1.5 + i*3.8)
    # Avatar circle
    avatar_shape = add_rect(s, x+Inches(1.1), Inches(2.0), Inches(0.9), Inches(0.9),
                            fill_color=ACCENT, border_color=None, radius=Pt(6))
    tf = avatar_shape.text_frame
    add_para(tf, avatar, Pt(24), align=PP_ALIGN.CENTER)
    add_text_box(s, x, Inches(3.05), Inches(3.0), Inches(0.3),
                 name, Pt(15), TEXT_1, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s, x, Inches(3.35), Inches(3.0), Inches(0.25),
                 role, Pt(11), ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, x, Inches(3.65), Inches(3.0), Inches(0.5),
                 desc, Pt(11), TEXT_2, align=PP_ALIGN.CENTER)

add_text_box(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.3),
             "GitHub: github.com/sleepycat583/ai-outfit-advisor", Pt(13), ACCENT, align=PP_ALIGN.CENTER, font_name='Consolas')
add_text_box(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.3),
             "技术栈: Python · LangChain · Chroma · Streamlit · DashScope · DuckDuckGo", Pt(12), TEXT_2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 16: THANKS
# ══════════════════════════════════════════
s = add_slide(); set_bg(s)
add_snum(s, 16)
add_text_box(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.5),
             "谢谢", Pt(96), ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.6),
             "小衣 · 你的私人穿搭主理人", Pt(22), TEXT_1, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.4),
             "基于大模型与 RAG 的大学生智能穿搭顾问", Pt(14), TEXT_3, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.4),
             "github.com/sleepycat583/ai-outfit-advisor", Pt(14), ACCENT, align=PP_ALIGN.CENTER, font_name='Consolas')
add_text_box(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
             "Q & A", Pt(28), ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.3),
             "欢迎评委老师提问", Pt(13), TEXT_2, align=PP_ALIGN.CENTER)

# ── Save ──
out_path = "c:/Users/user/Desktop/项目1/小衣_比赛答辩.pptx"
prs.save(out_path)
print(f"OK PPTX saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
